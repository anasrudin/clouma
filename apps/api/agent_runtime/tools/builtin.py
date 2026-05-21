"""
Builtin starter tools for Clouma agent runtime.

These tools are auto-registered when agent_runtime.tools is imported.

Security notes:
- read_file / run_python are unsandboxed — Phase 7+ concern.
- run_python executes arbitrary code in a subprocess with a timeout but no
  network or filesystem restrictions.
"""

from __future__ import annotations

import subprocess
import sys
import zoneinfo
from datetime import datetime
from pathlib import Path

import httpx

from . import register_tool


# ---------------------------------------------------------------------------
# web_search  (Tavily)
# ---------------------------------------------------------------------------


@register_tool(
    name="web_search",
    description=(
        "Search the web for a query and return a list of results, each with "
        "title, url, and snippet fields."
    ),
)
def web_search(query: str, max_results: int = 5) -> list[dict]:
    """
    Search the web via Tavily and return structured results.

    Args:
        query: The search query string.
        max_results: Maximum number of results to return (default 5).

    Returns:
        A list of dicts, each with keys: title, url, snippet.
    """
    from api.config import settings
    from tavily import TavilyClient

    if not settings.tavily_api_key:
        return [{"title": "web_search not configured", "url": "", "snippet": "Set TAVILY_API_KEY in .env to enable real web search."}]

    client = TavilyClient(
        api_key=settings.tavily_api_key,
        api_base_url=settings.tavily_base_url,
    )
    response = client.search(query, max_results=max_results)
    return [
        {
            "title": r.get("title", ""),
            "url": r.get("url", ""),
            "snippet": r.get("content", ""),
        }
        for r in response.get("results", [])
    ]


# ---------------------------------------------------------------------------
# http_get
# ---------------------------------------------------------------------------


@register_tool(
    name="http_get",
    description="Perform an HTTP GET request to the given URL and return the response body as text.",
)
def http_get(url: str) -> str:
    """
    Fetch the content of a URL via HTTP GET.

    Args:
        url: The URL to fetch.

    Returns:
        The response body as a string.
    """
    response = httpx.get(url, timeout=10)
    response.raise_for_status()
    return response.text


# ---------------------------------------------------------------------------
# read_file
# ---------------------------------------------------------------------------


@register_tool(
    name="read_file",
    description="Read the contents of a file at the given path and return it as text.",
)
def read_file(path: str) -> str:
    """
    Read a file from the filesystem.

    Args:
        path: Absolute or relative path to the file.

    Returns:
        The file contents as a UTF-8 string.

    Security note: This is unsandboxed — any path accessible to the process
    can be read.  Path allow-listing / sandboxing is a Phase 7+ concern.
    """
    return Path(path).read_text(encoding="utf-8")


# ---------------------------------------------------------------------------
# current_time
# ---------------------------------------------------------------------------


@register_tool(
    name="current_time",
    description=(
        "Return the current date and time in ISO 8601 format for the given "
        "timezone (default UTC)."
    ),
)
def current_time(tz: str = "UTC") -> str:
    """
    Get the current date and time in the specified timezone.

    Args:
        tz: IANA timezone name (e.g. 'UTC', 'America/New_York'). Default is 'UTC'.

    Returns:
        ISO 8601 datetime string (e.g. '2026-05-21T14:30:00+00:00').
    """
    return datetime.now(zoneinfo.ZoneInfo(tz)).isoformat()


# ---------------------------------------------------------------------------
# scrape_url
# ---------------------------------------------------------------------------

_SCRAPE_MAX_CHARS = 8_000
_REMOVE_TAGS = {"script", "style", "nav", "footer", "header", "aside", "noscript"}


@register_tool(
    name="scrape_url",
    description=(
        "Fetch a webpage and return its main text content, stripped of HTML tags, "
        "scripts, and navigation elements. Returns up to 8 000 characters."
    ),
)
def scrape_url(url: str) -> str:
    """
    Scrape a URL and return clean readable text.

    Args:
        url: The URL to scrape.

    Returns:
        Plain text extracted from the page body (up to 8 000 chars).
    """
    from bs4 import BeautifulSoup

    response = httpx.get(url, timeout=15, follow_redirects=True,
                         headers={"User-Agent": "Mozilla/5.0 (compatible; Clouma/1.0)"})
    response.raise_for_status()

    soup = BeautifulSoup(response.text, "html.parser")
    for tag in soup(list(_REMOVE_TAGS)):
        tag.decompose()

    text = soup.get_text(separator="\n", strip=True)
    # Collapse runs of blank lines to a single blank line
    lines = [l for l in text.splitlines() if l.strip()]
    cleaned = "\n".join(lines)
    return cleaned[:_SCRAPE_MAX_CHARS]


# ---------------------------------------------------------------------------
# run_python  (provider-based: subprocess | docker | e2b)
# ---------------------------------------------------------------------------

_SUBPROCESS_TIMEOUT = 10   # seconds
_DOCKER_TIMEOUT     = 30   # seconds — docker pull adds latency on first run
_E2B_TIMEOUT        = 60   # seconds — network round-trip


def _run_subprocess(code: str) -> str:
    """Unsandboxed fallback — runs in the same process space."""
    try:
        result = subprocess.run(
            [sys.executable, "-c", code],
            capture_output=True,
            text=True,
            timeout=_SUBPROCESS_TIMEOUT,
        )
    except subprocess.TimeoutExpired:
        return f"[error] Execution timed out after {_SUBPROCESS_TIMEOUT}s."

    if result.returncode != 0:
        return f"[error] Exit code {result.returncode}:\n{result.stderr.strip()}"
    return result.stdout


def _run_docker(code: str) -> str:
    """Run code inside a stripped-down Docker container with no network access."""
    try:
        result = subprocess.run(
            [
                "docker", "run", "--rm",
                "--network=none",           # no outbound network
                "--memory=128m",            # cap memory
                "--cpus=0.5",              # cap CPU
                "--security-opt=no-new-privileges",
                "python:3.11-slim",
                "python", "-c", code,
            ],
            capture_output=True,
            text=True,
            timeout=_DOCKER_TIMEOUT,
        )
    except FileNotFoundError:
        return "[error] Docker not found. Install Docker or set PYTHON_SANDBOX=subprocess."
    except subprocess.TimeoutExpired:
        return f"[error] Execution timed out after {_DOCKER_TIMEOUT}s."

    if result.returncode != 0:
        return f"[error] Exit code {result.returncode}:\n{result.stderr.strip()}"
    return result.stdout


def _run_e2b(code: str) -> str:
    """Run code in an E2B remote sandbox (https://e2b.dev)."""
    from api.config import settings

    if not settings.e2b_api_key:
        return "[error] E2B_API_KEY not set. Add it to .env or switch PYTHON_SANDBOX."

    from e2b_code_interpreter import Sandbox

    try:
        with Sandbox(api_key=settings.e2b_api_key, timeout=_E2B_TIMEOUT) as sbx:
            execution = sbx.run_code(code)
    except Exception as exc:
        return f"[error] E2B sandbox error: {exc}"

    if execution.error:
        return f"[error] {execution.error.name}: {execution.error.value}"
    return execution.text or ""


@register_tool(
    name="run_python",
    description=(
        "Execute a Python code snippet and return its stdout output. "
        "Provider is controlled by PYTHON_SANDBOX env var: "
        "'subprocess' (default, unsandboxed), 'docker' (local container, no network), "
        "'e2b' (remote sandbox). Use print() to produce output."
    ),
)
def run_python(code: str) -> str:
    """
    Execute a Python snippet and return stdout.

    Args:
        code: Valid Python source code to execute.

    Returns:
        stdout from the executed code, or an '[error] ...' message on failure.

    Provider is read from PYTHON_SANDBOX env var at call time:
      subprocess — unsandboxed, same process privileges (default)
      docker     — Docker container, --network=none, 128 MB RAM cap
      e2b        — E2B remote sandbox, requires E2B_API_KEY
    """
    from api.config import settings

    provider = settings.python_sandbox
    if provider == "docker":
        return _run_docker(code)
    if provider == "e2b":
        return _run_e2b(code)
    return _run_subprocess(code)


# ---------------------------------------------------------------------------
# write_file
# ---------------------------------------------------------------------------

@register_tool(
    name="write_file",
    description="Write text content to a file at the given path. mode='w' overwrites, mode='a' appends.",
)
def write_file(path: str, content: str, mode: str = "w") -> str:
    """
    Args:
        path: Absolute or relative path to write.
        content: Text content to write (UTF-8).
        mode: 'w' to overwrite (default), 'a' to append.
    Returns:
        The path that was written.
    """
    with open(path, mode, encoding="utf-8") as f:
        f.write(content)
    return path


# ---------------------------------------------------------------------------
# rss_fetch
# ---------------------------------------------------------------------------

@register_tool(
    name="rss_fetch",
    description="Fetch an RSS/Atom feed and return a list of entries with title, url, summary, published.",
)
def rss_fetch(url: str, max_items: int = 10) -> list[dict]:
    """
    Args:
        url: URL of the RSS or Atom feed.
        max_items: Maximum entries to return (default 10).
    Returns:
        List of dicts with keys: title, url, summary, published.
    """
    import feedparser
    feed = feedparser.parse(url)
    return [
        {
            "title": e.get("title", ""),
            "url": e.get("link", ""),
            "summary": e.get("summary", ""),
            "published": e.get("published", ""),
        }
        for e in feed.entries[:max_items]
    ]


# ---------------------------------------------------------------------------
# youtube_transcript
# ---------------------------------------------------------------------------

@register_tool(
    name="youtube_transcript",
    description="Fetch the transcript of a YouTube video and return it as plain text.",
)
def youtube_transcript(video_url: str, language: str = "en") -> str:
    """
    Args:
        video_url: Full YouTube URL (youtube.com/watch?v=... or youtu.be/...).
        language: BCP-47 language code for the transcript (default 'en').
    Returns:
        Full transcript as a single space-joined string, or '[error] ...' on failure.
    """
    import re
    from youtube_transcript_api import YouTubeTranscriptApi

    match = re.search(r"(?:v=|youtu\.be/)([A-Za-z0-9_-]{11})", video_url)
    if not match:
        return "[error] Could not extract video ID from URL."
    video_id = match.group(1)
    try:
        entries = YouTubeTranscriptApi.get_transcript(video_id, languages=[language])
        return " ".join(e["text"] for e in entries)
    except Exception as exc:
        return f"[error] {exc}"


# ---------------------------------------------------------------------------
# pdf_extract
# ---------------------------------------------------------------------------

_PDF_MAX_CHARS = 20_000


@register_tool(
    name="pdf_extract",
    description="Extract text from a PDF file (local path or URL). Returns up to 20 000 characters.",
)
def pdf_extract(path_or_url: str) -> str:
    """
    Args:
        path_or_url: Filesystem path or http(s) URL to a PDF.
    Returns:
        Extracted text (up to 20 000 chars), or '[error] ...' on failure.
    """
    import fitz  # pymupdf

    try:
        if path_or_url.startswith("http"):
            response = httpx.get(path_or_url, timeout=15, follow_redirects=True)
            response.raise_for_status()
            doc = fitz.open(stream=response.content, filetype="pdf")
        else:
            doc = fitz.open(path_or_url)
        text = "\n".join(page.get_text() for page in doc)
        doc.close()
        return text[:_PDF_MAX_CHARS]
    except Exception as exc:
        return f"[error] {exc}"


# ---------------------------------------------------------------------------
# pptx_generate
# ---------------------------------------------------------------------------

@register_tool(
    name="pptx_generate",
    description=(
        "Generate a PowerPoint presentation (.pptx). "
        "slides is a list of {title: str, bullets: list[str]}."
    ),
)
def pptx_generate(title: str, slides: list, output_path: str) -> str:
    """
    Args:
        title: Title of the presentation (shown on the title slide).
        slides: List of slide dicts: [{"title": "...", "bullets": ["...", ...]}].
        output_path: Where to save the .pptx file.
    Returns:
        output_path on success, '[error] ...' on failure.
    """
    try:
        from pptx import Presentation

        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title

        for slide_data in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_data.get("title", "")
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for bullet in slide_data.get("bullets", []):
                p = tf.add_paragraph()
                p.text = str(bullet)

        prs.save(output_path)
        return output_path
    except Exception as exc:
        return f"[error] {exc}"


# ---------------------------------------------------------------------------
# docx_generate
# ---------------------------------------------------------------------------

@register_tool(
    name="docx_generate",
    description=(
        "Generate a Word document (.docx). "
        "sections is a list of {heading: str, body: str}."
    ),
)
def docx_generate(sections: list, output_path: str) -> str:
    """
    Args:
        sections: List of section dicts: [{"heading": "...", "body": "..."}].
        output_path: Where to save the .docx file.
    Returns:
        output_path on success, '[error] ...' on failure.
    """
    try:
        from docx import Document

        doc = Document()
        for section in sections:
            if heading := section.get("heading"):
                doc.add_heading(str(heading), level=1)
            if body := section.get("body"):
                doc.add_paragraph(str(body))

        doc.save(output_path)
        return output_path
    except Exception as exc:
        return f"[error] {exc}"
